import io
import os
from io import BufferedReader
from pathlib import Path
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx
from loguru import logger
from typing import Union
import tempfile
import magic
from fastapi import HTTPException 
import tempfile
from models.models import Document, DocumentMetadata


async def get_document_from_file(
    file: Union[UploadFile, str], metadata: DocumentMetadata
) -> Document:
    extracted_text = ""
    if isinstance(file, UploadFile):
        file_content = await file.read()
        file = io.BytesIO(file_content)
        extracted_text = extract_text_from_file(file, file.content_type)
    elif isinstance(file, str) and Path(file).is_file():
        # When file_or_path is a filepath (string)
        extracted_text = extract_text_from_filepath(file)


    doc = Document(text=extracted_text, metadata=metadata)

    return doc

def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Récupère le contenu texte d'un fichier à partir de son chemin d'accès."""

    if mimetype is None:
        # Obtenez le type MIME du fichier en fonction de son extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Type de fichier non pris en charge")

    try:
        with open(filepath, "rb") as file:
            # Extrait le texte du fichier en utilisant la fonction extract_text_from_file
            extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        logger.error(e)
        raise e

    logger.info(f"Extracted text length: {len(extracted_text)}")
    
    return extracted_text

def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extrait le texte à partir d'un fichier PDF en utilisant PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Lit le texte à partir d'un fichier texte brut
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extrait le texte à partir d'un fichier docx en utilisant docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extrait le texte à partir d'un fichier CSV en utilisant le module csv
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extrait le texte à partir d'un fichier PPTX en utilisant python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Type de fichier non pris en charge
        raise ValueError("Type de fichier non pris en charge : {}".format(mimetype))

    logger.info(f"Extracted text length: {len(extracted_text)}")

    return extracted_text

# Extracte le texte d'un fichier en fonction de son type MIME
async def extract_text_from_form_file(file: UploadFile):
    """Renvoie le contenu textuel d'un fichier."""
    # Use python-magic to identify the file MIME type from its content
    # Read the file stream from the UploadFile object
    file_stream = await file.read()
    # Use the file_stream buffer for python-magic detection, not the UploadFile object
    mime_type = magic.from_buffer(file_stream, mime=True)
    logger.info(f"Identified MIME type using python-magic: {mime_type}")
    logger.info(f"file.file: {file.file}")
    logger.info("file: ", file)
    await file.seek(0)

    # Définit un chemin de fichier temporaire
    temp_file_path = "/tmp/temp_file"

    # Écrit le contenu du fichier dans un emplacement temporaire
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        # Extrait le texte en utilisant la fonction extract_text_from_filepath
        extracted_text = extract_text_from_filepath(temp_file_path, mime_type)
    except Exception as e:
        logger.error(e)
        os.remove(temp_file_path)
        raise e

    # Supprime le fichier de l'emplacement temporaire
    os.remove(temp_file_path)

    logger.info(f"Extracted text length: {len(extracted_text)}")

    return extracted_text
